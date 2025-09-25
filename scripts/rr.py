from kivy.app import App
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.textinput import TextInput
from kivy.uix.button import Button
from kivy.uix.label import Label
from kivy.uix.image import Image
from kivy.graphics import Color, Rectangle
from kivy.uix.filechooser import FileChooserIconView
from kivy.uix.popup import Popup
from kivy.clock import Clock
import os
import threading

# Librerías para documentos
import PyPDF2
import docx
from pptx import Presentation

# Intentamos importar SpeechRecognition solo para Windows/Linux
try:
    import speech_recognition as sr
    sr_available = True
except:
    sr_available = False

import pyttsx3
from plyer import tts as plyer_tts  # se deja, pero usamos pyttsx3 para control de stop


class FondoAzul(BoxLayout):
    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        with self.canvas.before:
            Color(0.89, 0.95, 0.99, 1)  # Azul muy suave
            self.rect = Rectangle(size=self.size, pos=self.pos)
        self.bind(size=self._update_rect, pos=self._update_rect)

    def _update_rect(self, *args):
        self.rect.size = self.size
        self.rect.pos = self.pos


class TTS_STTApp(App):
    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        self.current_language = "es-ES"
        self.text_input = None
        self.status_label = None

        # Control de hilos/stop
        self.tts_thread = None
        self._stop_event = threading.Event()
        self.engine_lock = threading.Lock()
        self.current_engine = None

    def build(self):
        main_layout = FondoAzul(orientation='vertical', padding=20, spacing=25)

        # Header
        header_layout = BoxLayout(orientation='horizontal', size_hint=(1, 0.2), spacing=15)
        logo_left = Image(source="inti-removebg-preview.png", size_hint=(None, 1), width=150)
        header_label = Label(
            text="INSTITUTO NACIONAL TÉCNICO INDUSTRIAL",
            color=(0.1, 0.2, 0.4, 1),
            font_size=32,
            bold=True,
            halign="center",
            valign="middle"
        )
        logo_right = Image(source="logo.png", size_hint=(None, 1), width=150)
        header_layout.add_widget(logo_left)
        header_layout.add_widget(header_label)
        header_layout.add_widget(logo_right)
        main_layout.add_widget(header_layout)

        # Status
        self.status_label = Label(
            text="Escribe o habla",
            size_hint=(1, 0.1),
            color=(0.1, 0.2, 0.4, 1),
            font_size=28,
            bold=True
        )
        main_layout.add_widget(self.status_label)

        # Caja de texto
        self.text_input = TextInput(
            hint_text=" Tu texto aparecerá aquí...",
            multiline=True,
            font_size=22,
            size_hint=(1, 0.5),
            background_normal="",
            background_color=(1, 1, 1, 1),
            foreground_color=(0, 0, 0, 1),
            padding=[20, 20],
            cursor_color=(0.25, 0.55, 0.95, 1)
        )
        main_layout.add_widget(self.text_input)

        # Botones de idioma
        lang_layout = BoxLayout(size_hint=(1, 0.15), spacing=20)
        btn_es = Button(text="Español", background_normal="", background_color=(0.20, 0.60, 0.90, 1), font_size=20, bold=True)
        btn_es.bind(on_press=lambda x: self.set_language("es-ES"))
        btn_en = Button(text="English", background_normal="", background_color=(0.20, 0.60, 0.90, 1), font_size=20, bold=True)
        btn_en.bind(on_press=lambda x: self.set_language("en-US"))
        lang_layout.add_widget(btn_es)
        lang_layout.add_widget(btn_en)
        main_layout.add_widget(lang_layout)

        # Botón documentos
        btn_doc = Button(
            text="Leer Documento",
            size_hint=(1, 0.15),
            background_normal="",
            background_color=(0.6, 0.4, 0.9, 1),
            font_size=22,
            bold=True
        )
        btn_doc.bind(on_press=self.open_filechooser)
        main_layout.add_widget(btn_doc)

        # Botones principales
        btn_layout = BoxLayout(size_hint=(1, 0.25), spacing=30)

        btn_tts = Button(text=" Texto a Voz ", size_hint=(0.33, 1),
                         background_normal="", background_color=(0.10, 0.75, 0.45, 1), font_size=22, bold=True)
        btn_tts.bind(on_press=self.text_to_speech)
        btn_layout.add_widget(btn_tts)

        btn_stt = Button(text=" Voz a Texto ", size_hint=(0.33, 1),
                         background_normal="", background_color=(0.10, 0.75, 0.45, 1), font_size=22, bold=True)
        btn_stt.bind(on_press=self.speech_to_text)
        btn_layout.add_widget(btn_stt)

        btn_stop = Button(text=" Detener Voz ", size_hint=(0.33, 1),
                          background_normal="", background_color=(0.9, 0.2, 0.2, 1), font_size=22, bold=True)
        btn_stop.bind(on_press=self.stop_speech)
        btn_layout.add_widget(btn_stop)

        main_layout.add_widget(btn_layout)

        return main_layout

    def set_language(self, lang_code):
        self.current_language = lang_code
        if lang_code == "es-ES":
            self.update_status("Idioma seleccionado: Español")
        else:
            self.update_status("Language selected: English")

    def update_status(self, text):
        Clock.schedule_once(lambda dt: setattr(self.status_label, 'text', text), 0)

    # === TTS mejorado con resaltado ===
    def text_to_speech(self, instance):
        text = self.text_input.text
        if not text.strip():
            self.update_status("❗ Escribe texto antes de reproducir")
            return

        self.stop_speech(None)
        self._stop_event.clear()
        self.tts_thread = threading.Thread(target=self._speak_text_thread, args=(text,), daemon=True)
        self.tts_thread.start()
        self.update_status("🔊 Reproduciendo... (pulse 'Detener Voz')")

    def _speak_text_thread(self, text):
        try:
            engine = pyttsx3.init()
        except Exception as e:
            print("Error al inicializar pyttsx3:", e)
            self.update_status("❌ Error inicializando TTS")
            return

        with self.engine_lock:
            self.current_engine = engine

        # Configurar voz
        try:
            voices = engine.getProperty('voices')
            for v in voices:
                name = (v.name or "").lower()
                if self.current_language.startswith("es") and "spanish" in name:
                    engine.setProperty('voice', v.id)
                    break
                if self.current_language.startswith("en") and "english" in name:
                    engine.setProperty('voice', v.id)
                    break
        except:
            pass

        engine.setProperty('rate', 150)
        engine.setProperty('volume', 1.0)

        # Dividir texto en párrafos
        paragraphs = [p.strip() for p in text.split("\n") if p.strip()]

        for para in paragraphs:
            if self._stop_event.is_set():
                break

            # Resaltar párrafo
            def highlight(dt, txt=para):
                self.text_input.select_text(0, 0)
                start = self.text_input.text.find(txt)
                if start != -1:
                    end = start + len(txt)
                    self.text_input.select_text(start, end)
                    self.text_input.cursor = (0, end)
                    self.text_input.scroll_y = 0
            Clock.schedule_once(highlight, 0)

            engine.say(para)
            engine.runAndWait()

        Clock.schedule_once(lambda dt: self.text_input.select_text(0, 0), 0)

        with self.engine_lock:
            self.current_engine = None

        if self._stop_event.is_set():
            self.update_status("⏹ Reproducción detenida")
        else:
            self.update_status("✅ Voz reproducida")

    def stop_speech(self, instance):
        self._stop_event.set()
        with self.engine_lock:
            if self.current_engine:
                try:
                    self.current_engine.stop()
                except Exception as e:
                    print("Error al detener engine:", e)
                finally:
                    self.current_engine = None
        self.update_status("⏹ Reproducción detenida")

    # === STT ===
    def speech_to_text(self, instance):
        if not sr_available:
            self.update_status("❌ STT no disponible")
            return
        threading.Thread(target=self._speech_to_text_thread, daemon=True).start()

    def _speech_to_text_thread(self):
        recognizer = sr.Recognizer()
        try:
            with sr.Microphone() as source:
                self.update_status("🎤 Escuchando...")
                audio = recognizer.listen(source, timeout=5, phrase_time_limit=20)
                text = recognizer.recognize_google(audio, language=self.current_language)
                Clock.schedule_once(lambda dt: setattr(self.text_input, 'text', text), 0)
                self.update_status(f"✅ Texto reconocido ({self.current_language})")
        except sr.WaitTimeoutError:
            self.update_status("⏱ No detecté voz")
        except sr.UnknownValueError:
            self.update_status("❌ No se entendió")
        except sr.RequestError:
            self.update_status("⚠ Error con Google")
        except Exception as e:
            print("Error en STT:", e)
            self.update_status("❌ Error en STT")

    # === Documentos ===
    def open_filechooser(self, instance):
        filechooser = FileChooserIconView(filters=["*.pdf", "*.docx", "*.pptx"], path=os.getcwd())
        popup = Popup(title="Selecciona un documento", content=filechooser, size_hint=(0.9, 0.9))
        filechooser.bind(on_submit=lambda chooser, selection, touch: self.load_document(selection, popup))
        popup.open()

    def load_document(self, selection, popup):
        if not selection:
            return
        file_path = selection[0]
        popup.dismiss()
        text = self.extract_text(file_path)
        if text:
            self.text_input.text = text
            self.update_status(f"📖 Documento cargado: {os.path.basename(file_path)}")
        else:
            self.update_status("❌ No se pudo leer el documento")

    def extract_text(self, file_path):
        text = ""
        if file_path.endswith(".pdf"):
            try:
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
            except Exception as e:
                print("Error leyendo PDF:", e)
                text = ""
        elif file_path.endswith(".docx"):
            try:
                doc = docx.Document(file_path)
                for para in doc.paragraphs:
                    text += para.text + "\n"
            except Exception as e:
                print("Error leyendo DOCX:", e)
                text = ""
        elif file_path.endswith(".pptx"):
            try:
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            except Exception as e:
                print("Error leyendo PPTX:", e)
                text = ""
        return text.strip()


if __name__ == "__main__":
    TTS_STTApp().run()


 